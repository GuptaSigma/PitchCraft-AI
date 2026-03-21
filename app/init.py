"""
PPTX Export Service
Author: GuptaSigma
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx. enum.text import PP_ALIGN
from io import BytesIO

class PPTXService:
    def __init__(self):
        self.prs = None
    
    def generate(self, presentation_data):
        """Generate PPTX file from presentation data"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Get slides
        slides_data = presentation_data.content. get('slides', [])
        
        for slide_data in slides_data:
            self._add_slide(slide_data)
        
        # Save to BytesIO
        output = BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output. getvalue()
    
    def _add_slide(self, slide_data):
        """Add a single slide"""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self. prs.slides.add_slide(blank_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame. text = slide_data.get('title', 'Untitled')
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Add content
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(9), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.text = slide_data.get('content', '')
        content_frame. paragraphs[0].font. size = Pt(18)
        content_frame.word_wrap = True

print("✅ PPTX Service loaded")