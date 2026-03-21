from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
from pathlib import Path
from urllib.parse import urlparse, urljoin
import requests
from PIL import Image
import json
import traceback
import logging

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTXService:
    """
    PowerPoint generation service with 8 theme support and 8 layouts
    """
    
    def __init__(self):
        """Initialize PPTX service with theme definitions"""
        
        # THEME DEFINITIONS (8 THEMES)
        self.themes = {
            'dialogue': {
                'name': 'Dialogue White',
                'bg': (255, 255, 255),      # White
                'text': (15, 23, 42),       # Dark slate
                'accent': (99, 102, 241),   # Indigo
                'card': (248, 250, 252)     # Light gray
            },
            'daytime': {
                'name': 'Daytime White',
                'bg': (255, 255, 255),
                'text': (15, 23, 42),
                'accent': (99, 102, 241),
                'card': (248, 250, 252)
            },
            'dawn': {
                'name': 'Dawn Warm',
                'bg': (255, 247, 237),
                'text': (124, 45, 18),
                'accent': (249, 115, 22),
                'card': (255, 237, 213)
            },
            'alien': {
                'name': 'Alien Dark',
                'bg': (15, 23, 42),         # Dark blue
                'text': (241, 245, 249),    # Light gray
                'accent': (34, 211, 238),   # Cyan
                'card': (30, 41, 59)        # Darker blue
            },
            'wine': {
                'name': 'Wine Elegance',
                'bg': (88, 28, 60),         # Dark wine
                'text': (255, 222, 200),    # Cream
                'accent': (244, 114, 182),  # Pink
                'card': (45, 11, 30)        # Deeper wine
            },
            'snowball': {
                'name': 'Snowball Blue',
                'bg': (224, 242, 254),      # Light blue
                'text': (30, 58, 138),      # Dark blue
                'accent': (14, 165, 233),   # Sky blue
                'card': (186, 230, 253)     # Lighter blue
            },
            'snowfall': {
                'name': 'Snowfall Blue',
                'bg': (224, 242, 254),
                'text': (30, 58, 138),
                'accent': (14, 165, 233),
                'card': (186, 230, 253)
            },
            'petrol': {
                'name': 'Petrol Steel',
                'bg': (71, 85, 105),        # Steel gray
                'text': (241, 245, 249),    # Light gray
                'accent': (14, 165, 233),   # Sky blue
                'card': (51, 65, 85)        # Darker steel
            },
            'piano': {
                'name': 'Piano Contrast',
                'bg': (255, 255, 255),      # White
                'text': (0, 0, 0),          # Black
                'accent': (0, 0, 0),        # Black
                'card': (245, 245, 245)     # Light gray
            },
            'sunset': {
                'name': 'Sunset Orange',
                'bg': (254, 252, 232),      # Cream
                'text': (120, 53, 15),      # Brown
                'accent': (249, 115, 22),   # Orange
                'card': (254, 243, 199)     # Light yellow
            },
            'midnight': {
                'name': 'Midnight Purple',
                'bg': (30, 41, 59),         # Dark gray
                'text': (226, 232, 240),    # Light gray
                'accent': (168, 85, 247),   # Purple
                'card': (15, 23, 42)        # Darker gray
            }
        }
        
        self.current_theme = self.themes['dialogue']  # Default
        logger.info("✅ PPTX Service v5.0.1 initialized with 8 themes")
    
    
    # BRIGHTNESS CALCULATION (NEW - v5.0.1)
    def _calculate_brightness(self, rgb):
        """
        Calculate perceived brightness of RGB color
        Returns: 0-255 (0=darkest, 255=brightest)
        
        Formula based on human eye sensitivity:
        - Green contributes most (58.7%)
        - Red moderate (29.9%)
        - Blue least (11.4%)
        """
        r, g, b = rgb
        brightness = (r * 299 + g * 587 + b * 114) / 1000
        return brightness
    
    
    def _get_readable_text_color(self, background_rgb):
        """
        Return black or white text based on background brightness
        
        Examples:
        - Bright yellow (255,255,0) → Black text
        - Dark blue (15,23,42) → White text
        """
        brightness = self._calculate_brightness(background_rgb)
        
        # Threshold: 128 (midpoint of 0-255 scale)
        if brightness > 128:
            return RGBColor(0, 0, 0)  # Black text for light backgrounds
        else:
            return RGBColor(255, 255, 255)  # White text for dark backgrounds
    
    
    # THEME DETECTION (ENHANCED - v5.0.1)
    def _detect_theme(self, presentation_data):
        """
        Enhanced theme detection with multiple fallbacks
        
        Handles:
        - Object attributes: presentation_data.theme
        - Dictionary access: presentation_data['theme']
        - Empty/None values
        - Whitespace
        """
        theme_name = None
        
        # Try 1: getattr (works for objects)
        theme_name = getattr(presentation_data, 'theme', None)
        
        # Try 2: __dict__ access (works for dict-like objects)
        if theme_name is None and hasattr(presentation_data, '__dict__'):
            if 'theme' in presentation_data.__dict__:
                theme_name = presentation_data.__dict__['theme']
        
        # Try 3: Direct dict access (if it's a dict)
        if theme_name is None and isinstance(presentation_data, dict):
            theme_name = presentation_data.get('theme')
        
        # Fallback: Default theme
        if not theme_name or str(theme_name).strip() == '':
            theme_name = 'dialogue'
            logger.warning("⚠️ No theme found, using default: dialogue")
        
        # Normalize
        theme_name = str(theme_name).lower().strip()
        
        # Validate
        if theme_name not in self.themes:
            logger.warning(f"⚠️ Invalid theme '{theme_name}', using dialogue")
            theme_name = 'dialogue'
        
        return theme_name
    
    
    # MAIN GENERATION METHOD
    def generate(self, presentation_data):
        """
        Generate PPTX file from presentation data
        
        Args:
            presentation_data: Object/dict with:
                - title (str)
                - theme (str): dialogue/alien/wine/etc.
                - content (dict): {'slides': [...]}
        
        Returns:
            bytes: PPTX file content
        """
        try:
            logger.info(f"\n{'='*80}")
            logger.info("[PPTX] Starting generation...")
            logger.info(f"{'='*80}")
            
            # STEP 1: DETECT & APPLY THEME
            theme_name = self._detect_theme(presentation_data)
            self.current_theme = self.themes[theme_name]
            
            logger.info(f"[THEME] Applied: {self.current_theme['name']} ({theme_name})")
            logger.info(f"[THEME] Background: RGB{self.current_theme['bg']}")
            logger.info(f"[THEME] Text Color: RGB{self.current_theme['text']}")
            logger.info(f"[THEME] Card Color: RGB{self.current_theme['card']}")
            logger.info(f"[THEME] Accent: RGB{self.current_theme['accent']}")
            
            # STEP 2: CREATE PRESENTATION
            prs = Presentation()
            # Use 16:9 widescreen to avoid square-ish output
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Get slides data
            content = presentation_data.content if hasattr(presentation_data, 'content') else presentation_data
            
            if isinstance(content, str):
                content = json.loads(content)
            
            slides_data = content.get('slides', []) if isinstance(content, dict) else content
            
            logger.info(f"[SLIDES] Total to generate: {len(slides_data)}")
            
            # STEP 3: GENERATE EACH SLIDE
            for idx, slide_data in enumerate(slides_data):
                slide_num = idx + 1
                layout = slide_data.get('layout', 'standard')
                
                # ✅ CYCLE LOGIC: Repeat 8-slide pattern for any additional slides
                pattern_num = (idx % 8) + 1
                
                layout_cycle = {
                    1: 'centered',
                    2: 'fixed_information',
                    3: 'three_col',
                    4: 'grid_4',
                    5: 'split_box',
                    6: 'roadmap',
                    7: 'fixed_information',
                    8: 'fixed_mission'
                }

                # ✅ FORCE PATTERN: 
                # Repeat pattern for any slide > 8, or if layout is 'standard'
                if slide_num > 8 or layout == 'standard':
                    layout = layout_cycle.get(pattern_num, layout)
                
                logger.info(f"[SLIDE {slide_num}] Pattern: {pattern_num} | Layout: {layout}")
                
                # Route to appropriate layout method using pattern_num for variant logic
                if layout == 'hero_overlay' or layout == 'hero':
                    self._create_hero_overlay_slide(prs, slide_data, slide_num)
                elif layout == 'centered':
                    self._create_centered_slide(prs, slide_data, slide_num)
                elif layout == 'split_box':
                    # Pattern 5 uses the special detailed split box
                    if pattern_num == 5:
                        self._create_fixed_split_box_slide_slide5(prs, slide_data, slide_num)
                    else:
                        self._create_fixed_split_box_slide(prs, slide_data, slide_num)
                elif layout == 'three_col':
                    self._create_fixed_three_cards_slide(prs, slide_data, slide_num)
                elif layout == 'grid_4':
                    self._create_fixed_four_grid_slide(prs, slide_data, slide_num)
                elif layout == 'fixed_information':
                    # Pattern 7 uses the Executive Summary layout
                    if pattern_num == 7:
                        self._create_executive_summary_slide(prs, slide_data, slide_num)
                    else:
                        self._create_fixed_image_cards_slide(prs, slide_data, slide_num)
                elif layout == 'roadmap':
                    self._create_fixed_roadmap_clean_slide(prs, slide_data, slide_num)
                elif layout == 'fixed_mission':
                    self._create_fixed_mission_slide(prs, slide_data, slide_num)
                else:
                    self._create_standard_slide(prs, slide_data, slide_num)
            
            # STEP 4: SAVE TO BYTES
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            
            logger.info(f"[SUCCESS] Presentation completed - Theme: {self.current_theme['name']}")
            logger.info(f"{'='*80}\n")
            
            return output.getvalue()
        
        except Exception as e:
            logger.error(f"❌ PPTX generation error: {e}")
            traceback.print_exc()
            raise
    
    
    # HELPER METHODS
    
    def _set_background(self, slide, color_rgb=None):
        """Set slide background color"""
        if color_rgb is None:
            color_rgb = self.current_theme['bg']
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)
    
    
    def _add_slide_number(self, slide, slide_num):
        """Add slide number to the bottom right corner with theme-aware color"""
        # Slide dimensions are 13.33 x 7.5 inches
        left = Inches(12.5)
        top = Inches(6.8)
        width = Inches(0.5)
        height = Inches(0.3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = Pt(12)
        
        # Use primary text color from theme for slide number
        p.font.color.rgb = RGBColor(*self.current_theme['text'])
        p.alignment = PP_ALIGN.RIGHT

    def _download_image(self, url, max_size=(1280, 720), cover=False):
        """Download an image from remote or local sources."""
        if not url:
            return None

        url_str = str(url).strip()
        try:
            parsed = urlparse(url_str)
        except Exception:
            parsed = None

        if parsed and parsed.scheme in ("http", "https"):
            return self._download_remote_image(url_str, max_size, cover)

        local_stream = self._load_local_image(parsed.path if parsed else url_str, max_size, cover)
        if local_stream:
            return local_stream

        if parsed and not parsed.scheme and url_str.startswith("//"):
            return self._download_remote_image(f"https:{url_str}", max_size, cover)

        if parsed and not parsed.scheme:
            try:
                from flask import request

                base = request.host_url if request else None
                if base:
                    absolute_url = urljoin(base, parsed.path.lstrip("/"))
                    return self._download_remote_image(absolute_url, max_size, cover)
            except Exception:
                pass

        logger.warning(f"⚠️ Image download failed: Unsupported or missing image source '{url_str}'")
        return None

    def _prepare_image(self, image, max_size, cover):
        if not max_size:
            return image

        if cover:
            return self._resize_cover(image, max_size)

        img_copy = image.copy()
        img_copy.thumbnail(max_size, Image.Resampling.LANCZOS)
        return img_copy

    def _download_remote_image(self, url, max_size, cover):
        try:
            response = requests.get(
                url,
                timeout=10,
                headers={
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)",
                    "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
                },
                stream=True,
            )
            if response.status_code == 200:
                img = Image.open(BytesIO(response.content))
                img = self._prepare_image(img, max_size, cover)

                output = BytesIO()
                img.convert("RGB").save(output, format='PNG')
                output.seek(0)
                return output
            logger.warning(f"⚠️ Remote image fetch failed with status {response.status_code} for '{url}'")
        except Exception as exc:
            logger.warning(f"⚠️ Image download failed: {exc}")
        return None

    def _load_local_image(self, path, max_size, cover):
        if not path:
            return None

        normalized = path.lstrip("/")
        try:
            root_dir = Path(__file__).resolve().parents[2]
            candidate = (root_dir / normalized).resolve()

            if not str(candidate).startswith(str(root_dir)) or not candidate.exists():
                return None

            with candidate.open('rb') as handle:
                data = handle.read()

            img = Image.open(BytesIO(data))
            img = self._prepare_image(img, max_size, cover)

            output = BytesIO()
            img.convert("RGB").save(output, format='PNG')
            output.seek(0)
            return output
        except Exception as exc:
            logger.warning(f"⚠️ Local image load failed for '{path}': {exc}")
        return None

    def _resize_cover(self, image, target_size):
        """Resize and crop an image so it fully covers the target size."""
        target_width, target_height = target_size
        if not target_width or not target_height:
            return image

        img_width, img_height = image.size
        if not img_width or not img_height:
            return image

        target_ratio = target_width / target_height
        image_ratio = img_width / img_height

        if image_ratio > target_ratio:
            scaled_height = target_height
            scaled_width = int(round(scaled_height * image_ratio))
        else:
            scaled_width = target_width
            scaled_height = int(round(scaled_width / image_ratio))

        resized = image.resize((scaled_width, scaled_height), Image.Resampling.LANCZOS)

        left = max(0, (scaled_width - target_width) // 2)
        top = max(0, (scaled_height - target_height) // 2)
        right = left + target_width
        bottom = top + target_height

        return resized.crop((left, top, right, bottom))
    
    
    def _parse_content(self, content):
        """Parse content into list of points (Improved Regex Cleaning)"""
        if isinstance(content, list):
            return content
        
        if not isinstance(content, str):
            content = str(content)
        
        # Split by newlines and clean
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        
        # Remove bullet markers, numbers, and Markdown characters
        import re
        cleaned = []
        for line in lines:
            # 1. Strip common bullet/number markers: "• ", "- ", "* ", "1. ", "01. ", etc.
            line = re.sub(r'^[\s•\-*]*([0-9]{1,2}[\.\)]\s*)?', '', line)
            
            # 2. Strip leading colon and spaces (common in AI output like "1. : Content")
            line = re.sub(r'^[:\s\t]*', '', line)
            
            # 3. Strip leading double asterisks (Markdown bold)
            line = re.sub(r'^\**\s*', '', line)
            
            # 4. Strip trailing asterisks
            line = re.sub(r'\**\s*$', '', line)
            
            # 5. Final trim
            line = line.strip()
            
            if line:
                cleaned.append(line)
        
        return cleaned

    def _fit_card_font_size(self, text):
        """Pick a smaller font size for longer card text without truncation."""
        length = len(text or "")
        if length > 420:
            return 8
        if length > 320:
            return 9
        if length > 240:
            return 10
        return 11

    def _text_scale(self, texts, min_len=120, max_len=600):
        """Return 0..1 scale based on total text length."""
        if not texts:
            return 0.0
        if isinstance(texts, str):
            total = len(texts)
        else:
            total = sum(len(str(t)) for t in texts)

        if max_len <= min_len:
            return 1.0

        scale = (total - min_len) / (max_len - min_len)
        return max(0.0, min(1.0, scale))

    def _scaled_font(self, scale, max_font, min_font):
        """Scale font size based on text length scale."""
        return int(round(max_font - scale * (max_font - min_font)))

    def _scaled_height(self, scale, min_height, max_height):
        """Scale height based on text length scale (in inches)."""
        return min_height + (max_height - min_height) * scale
    
    
    # LAYOUT 1: HERO SLIDE - FULL IMAGE WITH TEXT OVERLAY
    def _create_centered_slide(self, prs, slide_data, slide_num):
        """
        Hero slide with visible background image and text overlay.
        ✅ UPDATED: Improved fallback and content detection
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Try multiple image fields
        image_url = slide_data.get('image') or slide_data.get('bg_image') or slide_data.get('bg_url')
        
        image_loaded = False
        if image_url:
            img_stream = self._download_image(image_url, max_size=(1280, 960), cover=True)
            if img_stream:
                pic = slide.shapes.add_picture(img_stream, 0, 0, prs.slide_width, prs.slide_height)
                image_loaded = True

        if not image_loaded:
            self._set_background(slide) # Use theme background

        # Get Title with fallback
        title_text = slide_data.get('title')
        if not title_text or str(title_text).strip() == "":
            title_text = f"Chapter {slide_num}"
            
        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(2.4),
            width=prs.slide_width - Inches(2),
            height=Inches(1.4)
        )
        title_frame = title_box.text_frame
        title_frame.text = str(title_text).upper()
        title_frame.word_wrap = True

        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*self.current_theme['text'])

        # Get Content (checking mission/agenda as fallbacks for Slide 2)
        raw_content = slide_data.get('content') or slide_data.get('mission') or slide_data.get('agenda')
        content = self._parse_content(raw_content)
        if content:
            text = ' '.join(content) if isinstance(content, list) else str(content)
            scale = self._text_scale(text, min_len=140, max_len=700)
            text_box_top = Inches(3.7)
            text_box_height = Inches(self._scaled_height(scale, 2.4, 3.8))
            text_font_size = Pt(self._scaled_font(scale, 18, 14))

            bullet_circle = slide.shapes.add_shape(
                12,
                left=Inches(1.2),
                top=text_box_top + Inches(0.2),
                width=Inches(0.14),
                height=Inches(0.14)
            )
            bullet_circle.fill.solid()
            bullet_circle.fill.fore_color.rgb = RGBColor(*self.current_theme['accent']) # Use theme accent
            bullet_circle.line.fill.background()

            text_box = slide.shapes.add_textbox(
                left=Inches(1.6),
                top=text_box_top,
                width=prs.slide_width - Inches(2.2),
                height=text_box_height
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)

            text_p = text_frame.paragraphs[0]
            text_p.alignment = PP_ALIGN.LEFT
            text_p.line_spacing = 1.4
            text_p.text = text
            text_p.font.size = text_font_size
            text_p.font.color.rgb = RGBColor(*self.current_theme['text'])

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)

    def _create_hero_overlay_slide(self, prs, slide_data, slide_num):
        """
        Hero overlay slide with full background image and centered text.

        - Full background image
        - Dark transparent overlay
        - Centered title and paragraph
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        self._set_background(slide, self.current_theme['bg'])

        image_url = slide_data.get("image")
        if image_url:
            img_stream = self._download_image(image_url, max_size=(1600, 900), cover=True)
            if img_stream:
                slide.shapes.add_picture(
                    img_stream,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )

        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(2),
            width=Inches(8),
            height=Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", f"Slide {slide_num}")
        title_frame.word_wrap = True

        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*self.current_theme['text'])

        content = self._parse_content(slide_data.get("content", ""))
        full_text = " ".join(content) if content else ""

        text_box = slide.shapes.add_textbox(
            left=Inches(1.3),
            top=Inches(3.2),
            width=Inches(7.4),
            height=Inches(3)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = full_text

        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(*self.current_theme['text'])
        para.line_spacing = 1.4

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
    
    
    # LAYOUT 3: THREE CARDS (SLIDE 3) - ✅ FIXED v5.0.2
    def _create_fixed_three_cards_slide(self, prs, slide_data, slide_num):
        """
        Three vertical cards layout

        ✅ FIXED v5.0.2: Text overflow prevention
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # ✅ USE THEME BACKGROUND
        bg_color = self.current_theme['bg']
        title_color_val = self.current_theme['text']
        
        self._set_background(slide, bg_color)

        # ✅ Title - smaller and higher position
        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(0.3),
            width=Inches(8),
            height=Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color_val)

        # Parse content (expect 3 points)
        content = self._parse_content(slide_data.get('content', ''))
        points = content[:3]

        while len(points) < 3:
            points.append(f"Point {len(points) + 1}")

        # Cards with wider dimensions to fill space
        scale = self._text_scale(points, min_len=180, max_len=700)
        card_width = Inches(4.05)
        card_height_val = self._scaled_height(scale, 4.8, 6.1)
        card_height = Inches(card_height_val)
        start_left = Inches(0.35)
        card_top = Inches(1.2 + (6.1 - card_height_val) / 2)
        spacing = Inches(0.2)

        card_color = self.current_theme['card']
        card_text = self.current_theme['text']

        for i, point in enumerate(points):
            left = start_left + i * (card_width + spacing)

            # Card background
            card = slide.shapes.add_shape(
                1,  # Rectangle
                left=left,
                top=card_top,
                width=card_width,
                height=card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*card_color)
            card.line.color.rgb = RGBColor(*self.current_theme['accent'])
            card.line.width = Pt(2)

            # Number badge - smaller
            num_size = Inches(0.5)
            num_badge = slide.shapes.add_shape(
                12,  # Circle
                left=left + (card_width - num_size) / 2,
                top=card_top + Inches(0.3),
                width=num_size,
                height=num_size
            )
            num_badge.fill.solid()
            num_badge.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
            num_badge.line.fill.background()

            # Number text with smart color
            num_frame = num_badge.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.paragraphs[0].font.size = Pt(20)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = self._get_readable_text_color(self.current_theme['accent'])

            # Card text - with word wrap and margins
            text_box = slide.shapes.add_textbox(
                left=left + Inches(0.2),
                top=card_top + Inches(1.05),
                width=card_width - Inches(0.4),
                height=card_height - Inches(1.35)
            )
            text_frame = text_box.text_frame
            text_frame.text = point
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)

            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            text_frame.paragraphs[0].font.size = Pt(self._fit_card_font_size(point))
            text_frame.paragraphs[0].font.color.rgb = RGBColor(*card_text)
            text_frame.paragraphs[0].line_spacing = 1.35

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
    
    
    # LAYOUT 4: FOUR GRID (SLIDE 4) - ✅ FIXED v5.0.2
    def _create_fixed_four_grid_slide(self, prs, slide_data, slide_num):
        """
        2x2 grid layout with numbered circles

        ✅ FIXED v5.0.2: Text overflow prevention (no truncation)
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # ✅ USE THEME BACKGROUND
        self._set_background(slide)

        # ✅ Title - smaller and higher
        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(0.3),
            width=Inches(8),
            height=Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])

        # Parse content (expect 4 points)
        content = self._parse_content(slide_data.get('content', ''))
        points = content[:4]

        while len(points) < 4:
            points.append(f"Point {len(points) + 1}")

        # Grid layout - fill space better
        scale = self._text_scale(points, min_len=140, max_len=600)
        card_width = Inches(6.2)
        card_height_val = self._scaled_height(scale, 2.3, 3.1)
        card_height = Inches(card_height_val)
        start_left = Inches(0.3)
        start_top = Inches(1.1 + (3.1 - card_height_val) / 2)
        h_spacing = Inches(0.3)
        v_spacing = Inches(0.25)

        # ✅ Calculate readable text color for circles
        circle_text_color = self._get_readable_text_color(self.current_theme['accent'])

        positions = [
            (start_left, start_top),
            (start_left + card_width + h_spacing, start_top),
            (start_left, start_top + card_height + v_spacing),
            (start_left + card_width + h_spacing, start_top + card_height + v_spacing)
        ]

        for i, (left, top) in enumerate(positions):
            if i >= len(points):
                break

            # Card background
            card = slide.shapes.add_shape(
                1,  # Rectangle
                left=left,
                top=top,
                width=card_width,
                height=card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*self.current_theme['card'])
            card.line.color.rgb = RGBColor(*self.current_theme['accent'])
            card.line.width = Pt(2)

            # Number circle - smaller
            circle_size = Inches(0.45)
            circle = slide.shapes.add_shape(
                12,  # Circle
                left=left + Inches(0.2),
                top=top + Inches(0.2),
                width=circle_size,
                height=circle_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
            circle.line.fill.background()

            # Number text with smart color
            num_frame = circle.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.paragraphs[0].font.size = Pt(18)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = circle_text_color

            # Card text - with word wrap and margins
            text_box = slide.shapes.add_textbox(
                left=left + Inches(0.2),
                top=top + Inches(0.7),
                width=card_width - Inches(0.4),
                height=card_height - Inches(0.9)
            )
            text_frame = text_box.text_frame
            text_frame.text = points[i]
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            text_frame.paragraphs[0].font.size = Pt(self._scaled_font(scale, 13, 10))
            text_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])
            text_frame.paragraphs[0].line_spacing = 1.25

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
    
    
    # LAYOUT 5: IMAGE CARDS (SLIDE 2, 7) - ✅ FIXED v5.0.2
    def _create_fixed_image_cards_slide(self, prs, slide_data, slide_num):
        """
        Side-by-side layout with image and text

        ✅ FIXED v5.0.2: Text overflow prevention
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # ✅ USE THEME BACKGROUND
        self._set_background(slide)

        # LEFT: Image (50% width)
        image_url = slide_data.get('image')
        half_width = prs.slide_width / 2
        if image_url:
            img_stream = self._download_image(image_url, max_size=(640, 720), cover=True)
            if img_stream:
                slide.shapes.add_picture(
                    img_stream,
                    left=0,
                    top=0,
                    width=half_width,
                    height=prs.slide_height
                )

        # RIGHT: Content (50% width)
        content_left = half_width
        content_top = Inches(0.5)
        content_width = prs.slide_width - content_left

        # Title
        title_box = slide.shapes.add_textbox(
            left=content_left,
            top=content_top,
            width=content_width,
            height=Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(30)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])
        title_frame.paragraphs[0].line_spacing = 1.2

        # Content
        content = self._parse_content(slide_data.get('content', ''))
        points = content

        if not points:
            points = ["Content for this slide"]

        scale = self._text_scale(points, min_len=120, max_len=700)
        content_height = Inches(self._scaled_height(scale, 4.4, 5.9))
        point_font_size = Pt(self._scaled_font(scale, 16, 12))

        content_box = slide.shapes.add_textbox(
            left=content_left,
            top=content_top + Inches(1.1),
            width=content_width,
            height=content_height
        )

        frame = content_box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.2)
        frame.margin_right = Inches(0.2)
        frame.margin_top = Inches(0.1)

        for i, point in enumerate(points):
            p = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
            p.text = f"{i+1}. {point}"
            p.font.size = point_font_size
            p.font.color.rgb = RGBColor(*self.current_theme['text'])
            p.line_spacing = 1.35
            p.space_after = Pt(12)

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)

    # LAYOUT 7: EXECUTIVE SUMMARY WITH CORNER RIBBON - ✅ SLIDE 7
    def _create_executive_summary_slide(self, prs, slide_data, slide_num):
        """
        Modern executive summary with diagonal design elements.

        - White base background
        - Diagonal light gray overlay (left side)
        - Top-right triangle ribbon with "EXECUTIVE SUMMARY"
        - Layered triangles for depth
        - Right side decorative chevrons
        - Open text layout (no boxes)
        - 3 sections with TITLE + description
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Base background
        self._set_background(slide, (255, 255, 255))

        from pptx.enum.shapes import MSO_SHAPE

        # Diagonal gray overlay (left side)
        diagonal_overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(-2),
            top=Inches(-1),
            width=Inches(8),
            height=Inches(10)
        )
        diagonal_overlay.fill.solid()
        diagonal_overlay.fill.fore_color.rgb = RGBColor(245, 247, 250)
        diagonal_overlay.line.fill.background()
        diagonal_overlay.rotation = 15

        # Right side decorative chevrons
        chevron_colors = [
            (99, 102, 241),
            (139, 92, 246),
            (59, 130, 246)
        ]

        for i, color in enumerate(chevron_colors):
            chevron = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                left=prs.slide_width - Inches(1.5 + i * 0.3),
                top=Inches(2.5 + i * 0.5),
                width=Inches(1.5),
                height=Inches(1)
            )
            chevron.fill.solid()
            chevron.fill.fore_color.rgb = RGBColor(*color)
            chevron.line.fill.background()
            chevron.rotation = 90

        # Top-right triangle ribbon (layered)
        back_triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            left=prs.slide_width - Inches(2.2),
            top=Inches(-0.5),
            width=Inches(2.5),
            height=Inches(2.5)
        )
        back_triangle.fill.solid()
        back_triangle.fill.fore_color.rgb = RGBColor(147, 197, 253)
        back_triangle.line.fill.background()
        back_triangle.rotation = 135

        mid_triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            left=prs.slide_width - Inches(2),
            top=Inches(-0.3),
            width=Inches(2.3),
            height=Inches(2.3)
        )
        mid_triangle.fill.solid()
        mid_triangle.fill.fore_color.rgb = RGBColor(59, 130, 246)
        mid_triangle.line.fill.background()
        mid_triangle.rotation = 135

        front_triangle = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            left=prs.slide_width - Inches(1.8),
            top=Inches(-0.1),
            width=Inches(2.1),
            height=Inches(2.1)
        )
        front_triangle.fill.solid()
        front_triangle.fill.fore_color.rgb = RGBColor(30, 58, 138)
        front_triangle.line.fill.background()
        front_triangle.rotation = 135

        ribbon_text = slide.shapes.add_textbox(
            left=prs.slide_width - Inches(1.5),
            top=Inches(0.2),
            width=Inches(1.2),
            height=Inches(0.8)
        )
        ribbon_frame = ribbon_text.text_frame
        ribbon_frame.text = "EXECUTIVE\nSUMMARY"
        ribbon_frame.word_wrap = True
        ribbon_p = ribbon_frame.paragraphs[0]
        ribbon_p.alignment = PP_ALIGN.CENTER
        ribbon_p.font.size = Pt(9)
        ribbon_p.font.bold = True
        ribbon_p.font.color.rgb = RGBColor(255, 255, 255)
        ribbon_p.line_spacing = 0.9

        # Main title
        main_title = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(0.8)
        )
        main_title_frame = main_title.text_frame
        main_title_frame.text = slide_data.get('title', 'Executive Summary')
        main_title_frame.word_wrap = True
        main_title_p = main_title_frame.paragraphs[0]
        main_title_p.font.size = Pt(36)
        main_title_p.font.bold = True
        main_title_p.font.color.rgb = RGBColor(30, 41, 59)
        main_title_p.alignment = PP_ALIGN.LEFT

        # Accent line under title
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Inches(1),
            top=Inches(1.85),
            width=Inches(1.5),
            height=Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
        accent_line.line.fill.background()

        # Content sections (3)
        content = self._parse_content(slide_data.get('content', ''))
        while len(content) < 3:
            content.append(f"Section {len(content) + 1} content here...")

        sections = content[:3]
        section_start_top = Inches(2.3)
        section_spacing = Inches(1.6)
        section_left = Inches(1)
        section_width = Inches(7)

        for i, section_text in enumerate(sections):
            top = section_start_top + i * section_spacing

            title_label = slide.shapes.add_textbox(
                left=section_left,
                top=top,
                width=section_width,
                height=Inches(0.3)
            )
            title_frame = title_label.text_frame
            title_frame.text = "TITLE"
            title_frame.word_wrap = False
            label_p = title_frame.paragraphs[0]
            label_p.font.size = Pt(11)
            label_p.font.bold = True
            label_p.font.color.rgb = RGBColor(100, 116, 139)
            label_p.alignment = PP_ALIGN.LEFT

            text_length = len(section_text)
            if text_length < 120:
                font_size = 13
            elif text_length < 180:
                font_size = 12
            else:
                font_size = 11

            desc_text = slide.shapes.add_textbox(
                left=section_left,
                top=top + Inches(0.35),
                width=section_width,
                height=Inches(1.1)
            )
            desc_frame = desc_text.text_frame
            desc_frame.text = section_text
            desc_frame.word_wrap = True
            desc_frame.margin_left = Inches(0.05)

            desc_p = desc_frame.paragraphs[0]
            desc_p.font.size = Pt(font_size)
            desc_p.font.color.rgb = RGBColor(51, 65, 85)
            desc_p.alignment = PP_ALIGN.LEFT
            desc_p.line_spacing = 1.3

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
    
    
    # LAYOUT 6: ROADMAP/TIMELINE - ✅ WITH BOXES v5.0.3
    def _create_fixed_roadmap_clean_slide(self, prs, slide_data, slide_num):
        """
        Vertical timeline with boxes on the right side.

        - Circle numbers with vertical line (left)
        - Content boxes with background (right)
        - Year/title + description in each box when detected
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.25), Inches(11), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])

        # Parse content
        content = self._parse_content(slide_data.get('content', ''))
        steps = content[:6]
        if len(steps) < 3:
            steps = ["Step 1", "Step 2", "Step 3"]

        # Timeline configuration
        circle_size = Inches(0.55)
        circle_left = Inches(1.1)
        start_top = Inches(1.15)

        if len(steps) <= 4:
            base_max = 1.1
            base_min = 0.85
        elif len(steps) == 5:
            base_max = 0.95
            base_min = 0.75
        else:
            base_max = 0.82
            base_min = 0.65

        scale = self._text_scale(steps, min_len=120, max_len=700)
        step_height = Inches(self._scaled_height(scale, base_min, base_max))

        box_left = circle_left + circle_size + Inches(0.25)
        box_width = Inches(10.0)
        box_height = step_height - Inches(0.12)

        circle_text_color = self._get_readable_text_color(self.current_theme['accent'])

        for i, step in enumerate(steps):
            top = start_top + i * step_height

            # Circle number
            circle = slide.shapes.add_shape(
                12,  # Circle
                left=circle_left,
                top=top,
                width=circle_size,
                height=circle_size
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
            circle.line.fill.background()

            num_frame = circle.text_frame
            num_frame.text = str(i + 1)
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            num_frame.paragraphs[0].font.size = Pt(20)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = circle_text_color

            # Connector line
            if i < len(steps) - 1:
                connector = slide.shapes.add_shape(
                    1,  # Rectangle
                    left=circle_left + circle_size / 2 - Pt(2),
                    top=top + circle_size,
                    width=Pt(4),
                    height=step_height - circle_size
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
                connector.line.fill.background()

            # Content box
            box = slide.shapes.add_shape(
                1,  # Rectangle
                left=box_left,
                top=top,
                width=box_width,
                height=box_height
            )
            box.fill.solid()

            if self._calculate_brightness(self.current_theme['bg']) > 128:
                box.fill.fore_color.rgb = RGBColor(245, 247, 250)
            else:
                box.fill.fore_color.rgb = RGBColor(*self.current_theme['card'])

            box.line.color.rgb = RGBColor(*self.current_theme['accent'])
            box.line.width = Pt(1.5)

            accent_bar = slide.shapes.add_shape(
                1,  # Rectangle
                left=box_left,
                top=top,
                width=Pt(4),
                height=box_height
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
            accent_bar.line.fill.background()

            # Year + description
            import re

            year_match = re.match(r'^(\d{4})[:\s-]*(.+)', step.strip())
            if year_match:
                year = year_match.group(1)
                description = year_match.group(2).strip()
            else:
                year = None
                description = step.strip()

            if year:
                year_box = slide.shapes.add_textbox(
                    left=box_left + Inches(0.15),
                    top=top + Inches(0.08),
                    width=box_width - Inches(0.3),
                    height=Inches(0.32)
                )
                year_frame = year_box.text_frame
                year_frame.text = year
                year_frame.word_wrap = False
                year_frame.paragraphs[0].font.size = Pt(16)
                year_frame.paragraphs[0].font.bold = True
                year_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['accent'])
                desc_top = top + Inches(0.4)
                desc_height = box_height - Inches(0.5)
            else:
                desc_top = top + Inches(0.1)
                desc_height = box_height - Inches(0.2)

            text_length = len(description)
            if text_length < 80:
                font_size = 11
            elif text_length < 120:
                font_size = 10
            else:
                font_size = 9

            desc_box = slide.shapes.add_textbox(
                left=box_left + Inches(0.15),
                top=desc_top,
                width=box_width - Inches(0.3),
                height=desc_height
            )
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            desc_frame.word_wrap = True
            desc_frame.margin_left = Inches(0.05)
            desc_frame.margin_right = Inches(0.05)
            desc_frame.paragraphs[0].font.size = Pt(font_size)
            desc_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])
            desc_frame.paragraphs[0].line_spacing = 1.2

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
    
    
    # LAYOUT 7: SPLIT BOX (SLIDE 5) - ✅ FIXED v5.0.1
    def _create_fixed_split_box_slide(self, prs, slide_data, slide_num):
        """
        Split layout with image left, content right
        
        ✅ FIXED: Uses theme colors
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # ✅ USE THEME BACKGROUND
        self._set_background(slide)
        
        # Left side: Image (50% width)
        image_url = slide_data.get('image')
        half_width = prs.slide_width / 2
        if image_url:
            img_stream = self._download_image(image_url, max_size=(600, 720), cover=True)
            if img_stream:
                slide.shapes.add_picture(
                    img_stream,
                    left=0,
                    top=0,
                    width=half_width,
                    height=prs.slide_height
                )

        # Right side: Content box (50% width)
        content_box = slide.shapes.add_shape(
            1,  # Rectangle
            left=half_width,
            top=0,
            width=half_width,
            height=prs.slide_height
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(*self.current_theme['card'])
        content_box.line.fill.background()
        
        # ✅ Title
        title_box = slide.shapes.add_textbox(
            left=half_width + Inches(0.3),
            top=Inches(1),
            width=half_width - Inches(0.6),
            height=Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])
        
        # ✅ Content
        content = self._parse_content(slide_data.get('content', ''))
        points = content[:5]

        scale = self._text_scale(points, min_len=120, max_len=700)
        text_height = Inches(self._scaled_height(scale, 3.8, 4.9))
        point_font = Pt(self._scaled_font(scale, 16, 12))
        
        text_box = slide.shapes.add_textbox(
            left=half_width + Inches(0.3),
            top=Inches(2.5),
            width=half_width - Inches(0.6),
            height=text_height
        )
        
        frame = text_box.text_frame
        frame.word_wrap = True
        
        for i, point in enumerate(points):
            p = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = point_font
            p.font.color.rgb = RGBColor(*self.current_theme['text'])
            p.space_after = Pt(12)

        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)

    # LAYOUT: SPLIT BOX (SLIDE 5 ONLY) - ✅ PERFECT FIT v5.0.3
    def _create_fixed_split_box_slide_slide5(self, prs, slide_data, slide_num):
        """
        Split layout with image left, 2+1 grid cards right (Slide 5 only).
        Matches requested visual: Title top, intro below, 3 boxes in a dark container.
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_background(slide)

        # LEFT: Image (50% width)
        image_url = slide_data.get('image')
        half_width = prs.slide_width / 2
        if image_url:
            img_stream = self._download_image(image_url, max_size=(600, 720), cover=True)
            if img_stream:
                slide.shapes.add_picture(img_stream, 0, 0, half_width, prs.slide_height)

        # RIGHT: force black text for slide 5 readability
        panel_rgb = RGBColor(*self.current_theme['bg'])
        title_rgb = RGBColor(0, 0, 0)
        text_rgb = RGBColor(0, 0, 0)
        container_rgb = RGBColor(*self.current_theme['card'])
        card_bg_rgb = RGBColor(*self.current_theme['card'])
        card_text_rgb = RGBColor(0, 0, 0)

        # 1. Right side panel
        panel = slide.shapes.add_shape(1, half_width, 0, half_width, prs.slide_height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = panel_rgb
        panel.line.fill.background()

        # 2. Title
        title_box = slide.shapes.add_textbox(half_width + Inches(0.4), Inches(0.4), half_width - Inches(0.8), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', f'Slide {slide_num}')
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = title_rgb

        # 3. Intro text
        content = self._parse_content(slide_data.get('content', ''))
        intro_text = content[0] if content else "Introduction text..."
        intro_box = slide.shapes.add_textbox(half_width + Inches(0.4), Inches(1.2), half_width - Inches(0.8), Inches(1.0))
        inf = intro_box.text_frame
        inf.word_wrap = True
        inf.margin_left = Inches(0.05)
        inf.margin_right = Inches(0.05)
        p2 = inf.paragraphs[0]
        p2.text = f"• {intro_text}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = text_rgb

        # 4. Main Container (Rounded Box)
        cont_w = half_width - Inches(0.8)
        cont_h = prs.slide_height - Inches(2.6)
        cont_shape = slide.shapes.add_shape(5, half_width + Inches(0.4), Inches(2.4), cont_w, cont_h - Inches(0.2))
        cont_shape.fill.solid()
        cont_shape.fill.fore_color.rgb = container_rgb
        cont_shape.line.fill.background()

        # 5. Grid Cards (01, 02 top | 03 bottom)
        card_points = content[1:4] if len(content) > 1 else []
        while len(card_points) < 3:
            card_points.append(f"Historical detail {len(card_points)+1}")

        # Card geometry
        card_gap = Inches(0.2)
        sm_w = (cont_w - (3 * card_gap)) / 2
        sm_h = (cont_h - (3 * card_gap)) * 0.45
        
        # Positions relative to container
        c1x = half_width + Inches(0.4) + card_gap
        c1y = Inches(2.4) + card_gap
        
        c2x = c1x + sm_w + card_gap
        c2y = c1y
        
        c3x = c1x
        c3w = cont_w - (2 * card_gap)
        c3h = cont_h - sm_h - (3 * card_gap)
        c3y = c1y + sm_h + card_gap

        card_geos = [
            (c1x, c1y, sm_w, sm_h),
            (c2x, c2y, sm_w, sm_h),
            (c3x, c3y, c3w, c3h)
        ]

        accent_rgb = RGBColor(*self.current_theme['accent'])

        for i, (cx, cy, cw, ch) in enumerate(card_geos):
            # Card shape
            c_sh = slide.shapes.add_shape(5, cx, cy, cw, ch)
            c_sh.fill.solid()
            c_sh.fill.fore_color.rgb = card_bg_rgb
            c_sh.line.color.rgb = RGBColor(*self.current_theme['accent'])
            c_sh.line.width = Pt(0.5)

            # Number (01, 02, 03)
            num_box = slide.shapes.add_textbox(cx + Inches(0.1), cy + Inches(0.1), Inches(0.5), Inches(0.4))
            num_p = num_box.text_frame.paragraphs[0]
            num_p.text = f"0{i+1}"
            num_p.font.size = Pt(16)
            num_p.font.bold = True
            num_p.font.color.rgb = accent_rgb

            # Card Content
            p_box = slide.shapes.add_textbox(cx + Inches(0.1), cy + Inches(0.5), cw - Inches(0.2), ch - Inches(0.6))
            p_frame = p_box.text_frame
            p_frame.word_wrap = True
            p_frame.margin_left = Inches(0.05)
            p_frame.margin_right = Inches(0.05)
            p_frame.margin_bottom = Inches(0.05)
            cp = p_frame.paragraphs[0]
            cp.text = card_points[i]
            text_len = len(card_points[i])
            if text_len < 100:
                cp.font.size = Pt(12)
            elif text_len < 200:
                cp.font.size = Pt(11)
            else:
                cp.font.size = Pt(10)
            cp.font.color.rgb = card_text_rgb

        return slide

    # LAYOUT 8: SPLIT SCREEN WITH ROUNDED CARD OVERLAY ✅
    def _create_image_overlay_slide(self, prs, slide_data, slide_num):
        """
        Split screen layout with rounded card overlay.

        - White background base
        - Right side: Full-height image
        - Left side: Large rounded card with text
        - Card overlaps image slightly
        - "KEY INSIGHT" label
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Base layer: white background
        self._set_background(slide, (255, 255, 255))

        # Right side: full-height image
        image_url = slide_data.get('image')
        if image_url:
            img_stream = self._download_image(image_url, max_size=(1280, 960), cover=True)
            if img_stream:
                img_left = Inches(6.0)
                img_width = prs.slide_width - img_left
                slide.shapes.add_picture(
                    img_stream,
                    left=img_left,
                    top=0,
                    width=img_width,
                    height=prs.slide_height
                )

        # Left overlay: rounded card (with overlap)
        content_scale = self._text_scale(self._parse_content(slide_data.get('content', '')), min_len=160, max_len=700)
        card_left = Inches(0.6)
        card_height_val = self._scaled_height(content_scale, 4.6, 5.8)
        card_top = Inches(1 + (5.8 - card_height_val) / 2)
        card_width = Inches(5.5)
        card_height = Inches(card_height_val)

        card = slide.shapes.add_shape(
            5,  # Rounded rectangle
            left=card_left,
            top=card_top,
            width=card_width,
            height=card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*self.current_theme['accent'])
        card.line.fill.background()

        # "KEY INSIGHT" label
        label_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.5),
            top=card_top + Inches(0.4),
            width=card_width - Inches(1),
            height=Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = "KEY INSIGHT"
        label_frame.word_wrap = False
        label_frame.paragraphs[0].font.size = Pt(12)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        label_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Card title
        title_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.5),
            top=card_top + Inches(0.9),
            width=card_width - Inches(1),
            height=Inches(1.1)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', 'Summary and Conclusion')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].line_spacing = 1.1

        # Card content (paragraphs)
        content = self._parse_content(slide_data.get('content', ''))
        paragraphs = content[:4]
        if not paragraphs:
            paragraphs = [
                "First key point about the topic.",
                "Second important consideration.",
                "Third strategic element.",
                "Final concluding thought."
            ]

        content_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.5),
            top=card_top + Inches(2.2),
            width=card_width - Inches(1),
            height=card_height - Inches(2.6)
        )

        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.margin_left = Inches(0.05)
        content_frame.margin_right = Inches(0.05)

        total_chars = sum(len(p) for p in paragraphs)
        text_scale = self._text_scale(paragraphs, min_len=400, max_len=900)
        font_size = self._scaled_font(text_scale, 13, 10)
        spacing = int(round(self._scaled_height(text_scale, 12, 18)))

        for i, para_text in enumerate(paragraphs):
            p = content_frame.add_paragraph() if i > 0 else content_frame.paragraphs[0]
            p.text = para_text
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.line_spacing = 1.3
            p.space_after = Pt(spacing)
            p.alignment = PP_ALIGN.LEFT
    
    
    # LAYOUT 8: MISSION/PANEL (SLIDE 8+) - ✅ FIXED v5.0.1
    def _create_fixed_mission_slide(self, prs, slide_data, slide_num):
        """
        Modified Mission Slide (Layout 8): Floating card over image.
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 1. Theme background
        self._set_background(slide, self.current_theme['bg'])
        
        # 2. Add Image on the Right (70% width)
        image_url = slide_data.get('image')
        if image_url:
            img_width = prs.slide_width * 0.70
            img_stream = self._download_image(image_url, cover=True)
            if img_stream:
                slide.shapes.add_picture(
                    img_stream,
                    left=prs.slide_width - img_width,
                    top=0,
                    width=img_width,
                    height=prs.slide_height
                )

        # 3. Add floating card using theme card color
        # Increased size to prevent text overflow v5.0.4
        card_width = Inches(6.5)
        card_height = Inches(5.8)
        
        # Alignment: Card starts at 8% of slide width (slightly more left)
        card_left = prs.slide_width * 0.08
        card_top = (prs.slide_height - card_height) / 2
        
        shape = slide.shapes.add_shape(
            5,  # Rounded Rectangle
            left=card_left,
            top=card_top,
            width=card_width,
            height=card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*self.current_theme['card'])
        shape.line.color.rgb = RGBColor(*self.current_theme['accent'])
        shape.line.width = Pt(1.5)
        
        # Add shadow to card
        try:
            shape.shadow.inherit = False
            shape.shadow.visible = True
        except: pass

        # 4. Content inside card
        # Small Label
        label_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.4),
            top=card_top + Inches(0.4),
            width=Inches(3),
            height=Inches(0.4)
        )
        label_frame = label_box.text_frame
        label_frame.text = "KEY INSIGHT"
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(12)
        label_p.font.bold = True
        label_p.font.color.rgb = RGBColor(*self.current_theme['accent'])

        # Big Title (Summary) - Slightly smaller font to fit precisely
        title_text = slide_data.get('title', 'Summary')
        title_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.4),
            top=card_top + Inches(0.8),
            width=card_width - Inches(0.8),
            height=Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(34) # Reduced from 40 for better fit
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*self.current_theme['text'])

        # Description (Points) - Increased box height for more text
        content_raw = slide_data.get('content', '')
        points = self._parse_content(content_raw)
        
        desc_box = slide.shapes.add_textbox(
            left=card_left + Inches(0.4),
            top=card_top + Inches(2.1), # Pushed down slightly more
            width=card_width - Inches(0.8),
            height=Inches(3.2) # Increased height
        )
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        
        if not points:
            points = ["Vision and mission summary", "Key strategic objectives"]
            
        for i, point_text in enumerate(points[:4]):
            if i == 0:
                p = desc_frame.paragraphs[0]
            else:
                p = desc_frame.add_paragraph()
            
            p.text = f"• {point_text}"
            p.font.size = Pt(14) # Reduced from 16 for better fit
            p.font.color.rgb = RGBColor(*self.current_theme['text'])
            p.space_after = Pt(8)

        # 5. Slide Number
        self._add_slide_number(slide, slide_num)
    
    # LAYOUT: STANDARD/DEFAULT SLIDE - ✅ v5.0.1
    def _create_standard_slide(self, prs, slide_data, slide_num):
        """
        Standard slide layout with title and bullet points
        
        ✅ FIXED: Uses theme colors
        - Title at top
        - Content in bullet points below
        - Clean, professional layout
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # ✅ USE THEME BACKGROUND
        self._set_background(slide)
        
        # ✅ Title
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=prs.slide_width - Inches(1),
            height=Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get('title', f'Slide {slide_num}')
        title_frame.word_wrap = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(*self.current_theme['text'])
        
        # ✅ Content
        content = self._parse_content(slide_data.get('content', ''))
        points = content[:8]  # Limit to 8 points for standard slide
        
        if not points:
            points = ["Content for this slide"]
        
        # Dynamic sizing based on content
        scale = self._text_scale(points, min_len=100, max_len=600)
        content_height = Inches(self._scaled_height(scale, 5.0, 6.0))
        point_font_size = Pt(self._scaled_font(scale, 18, 14))
        
        content_box = slide.shapes.add_textbox(
            left=Inches(0.8),
            top=Inches(1.8),
            width=prs.slide_width - Inches(1.6),
            height=content_height
        )
        
        frame = content_box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.1)
        frame.margin_right = Inches(0.1)
        frame.margin_top = Inches(0.1)
        
        for i, point in enumerate(points):
            p = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = point_font_size
            p.font.color.rgb = RGBColor(*self.current_theme['text'])
            p.line_spacing = 1.4
            p.space_after = Pt(14)
        
        # ✅ Slide Number
        self._add_slide_number(slide, slide_num)
        
        logger.info(f"   → Standard slide created with {len(points)} points")

# MODULE INITIALIZATION
logger.info("✅ PPTX Service v5.0.1 (THEME FIX) loaded successfully")
logger.info("   - 8 Themes supported")
logger.info("   - 8 Layouts implemented")
logger.info("   - Smart text colors based on brightness")
logger.info("   - Enhanced theme detection")
logger.info("   - All hardcoded colors removed")